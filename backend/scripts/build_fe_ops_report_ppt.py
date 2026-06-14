"""기능인인정제 서명 시스템 — 대표이사 보고용 운영설명서 PPT 생성."""
from __future__ import annotations

from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[2]
SHOT_DIR = ROOT / "docs" / "reports" / "functional-eval-e2e" / "screenshots" / "report"
OUT_PPT = ROOT / "docs" / "기능인인정제_서명시스템_운영설명서_보고용.pptx"

NAVY = RGBColor(0x1E, 0x3A, 0x5F)
BLUE = RGBColor(0x25, 0x63, 0xEB)
GRAY = RGBColor(0x64, 0x74, 0x8B)
DARK = RGBColor(0x0F, 0x17, 0x2A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TABLE_HEAD = RGBColor(0xE2, 0xE8, 0xF0)


def _set_slide_bg(slide, color: RGBColor = WHITE) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_title(slide, title: str, subtitle: str | None = None) -> None:
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.0), Inches(0.9))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = NAVY
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.6), Inches(1.05), Inches(12.0), Inches(0.45))
        stf = sub.text_frame
        stf.text = subtitle
        stf.paragraphs[0].font.size = Pt(14)
        stf.paragraphs[0].font.color.rgb = GRAY


def _add_bullets(slide, lines: list[str], top: float = 1.55, height: float = 1.2, width: float = 12.0) -> None:
    box = slide.shapes.add_textbox(Inches(0.6), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.color.rgb = DARK
        p.space_after = Pt(6)
        p.level = 0


def _add_image(slide, filename: str, left: float, top: float, width: float) -> bool:
    path = SHOT_DIR / filename
    if not path.is_file():
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.4))
        box.text_frame.text = f"[캡처 없음: {filename}]"
        box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xB9, 0x1C, 0x1C)
        return False
    slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width))
    return True


def _slide_screen(
    prs: Presentation,
    title: str,
    image: str | None,
    bullets: list[str],
    notes: list[str] | None = None,
    image2: str | None = None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _add_title(slide, title)
    _add_bullets(slide, bullets, top=1.45, height=1.35)
    y = 2.85
    if image:
        _add_image(slide, image, 0.6, y, 7.8)
    if image2:
        _add_image(slide, image2, 8.6, y, 4.0)
    if notes:
        _add_bullets(slide, notes, top=6.55, height=0.8, width=12.2)


def _slide_table(prs: Presentation, title: str, headers: list[str], rows: list[list[str]]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _add_title(slide, title)
    cols, row_count = len(headers), len(rows) + 1
    table_shape = slide.shapes.add_table(row_count, cols, Inches(0.6), Inches(1.6), Inches(12.0), Inches(0.45 * row_count))
    table = table_shape.table
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEAD
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(12)
            p.font.color.rgb = NAVY
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.color.rgb = DARK if val != "미구현" else RGBColor(0xB9, 0x1C, 0x1C)


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1 표지
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    t = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.2))
    tp = t.text_frame.paragraphs[0]
    tp.text = "기능인인정제 서명 시스템 운영 설명서"
    tp.font.size = Pt(36)
    tp.font.bold = True
    tp.font.color.rgb = NAVY
    tp.alignment = PP_ALIGN.CENTER
    s = slide.shapes.add_textbox(Inches(0.8), Inches(3.2), Inches(11.5), Inches(0.6))
    sp = s.text_frame.paragraphs[0]
    sp.text = "평가자 책임확인 · 전자서명 · 평가보고서 자동보관 체계"
    sp.font.size = Pt(18)
    sp.font.color.rgb = BLUE
    sp.alignment = PP_ALIGN.CENTER
    meta = slide.shapes.add_textbox(Inches(0.8), Inches(5.6), Inches(11.5), Inches(1.0))
    mtf = meta.text_frame
    mtf.text = f"작성일: {date.today():%Y-%m-%d}  |  작성: BESMA 개발팀  |  시스템: BESMA 기능인인정제(웹)"
    mtf.paragraphs[0].font.size = Pt(13)
    mtf.paragraphs[0].font.color.rgb = GRAY
    mtf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 2 구축 목적
    _slide_screen(
        prs,
        "1. 구축 목적",
        None,
        [
            "기존 기능인인정제(2-1 기능 · 2-2 안전) 평가 흐름에 전자서명·결재 단계를 접목",
            "평가자의 「평가 완료」 선언을 캔버스 서명·PDF로 증빙하고 서버에 자동 보관",
            "최초 로그인 시 평가책임 동의서 1회 서명(이후 재표시 없음)",
            "평가완료보고서 서명 후 해당 평가·근거(포상/제재) 수정 불가",
            "신규 입사·추가 출역 인원은 기존 서명본을 수정하지 않고 추가평가(batch)로 분리 처리",
        ],
    )

    # 3 전체 프로세스
    _slide_screen(
        prs,
        "2. 전체 운영 프로세스",
        None,
        [
            "로그인 → 최초 동의서 서명 → 근로자 2-1/2-2 평가 → 포상·제재 등록",
            "→ (11명↑ 현장) 팀장 평가완료보고서 서명 → 소장 팀장보고서 승인·서명",
            "→ 소장 현장 평가보고서 제출·서명 → 안전보건실 담당 검토·승인",
            "→ 안전보건실장 최종 승인·서명 → 대표이사 최종 승인·서명 → PDF 보고서 저장",
            "※ 10명 이하 현장: 팀장 단계 생략, 소장이 전원 평가·제출",
        ],
    )

    # 화면별
    # 평가 기준 (기능/품질 vs 안전)
    _slide_screen(
        prs,
        "3-0. 평가 기준 — 기능/품질 vs 안전",
        "04_evaluate_functional.png",
        [
            "2-1 기능/품질: 소장·팀장 평가 — 작업숙련도·품질·생산성·작업태도",
            "기본 B등급 · S등급 현장별 20% 이내 권장 · 초과 시 평가완료 서명 전 사유",
            "C등급: 작업능력 미흡·품질 불량·재작업·태도 불량 등 근거 있을 때만 (비율 강제 없음)",
            "2-2 안전: 감점형 절대평가 — 지적·위반 이력 없으면 감점하지 않음(만점/상위)",
            "안전 S 20% 제한 없음 · 안전 C 최소 비율 강제 없음",
        ],
        [
            "기능/품질 평가는 소장·팀장이 작업숙련도와 품질을 기준으로 평가하고, S등급 남발 방지를 위해 기본 B등급 및 S등급 권장 상한을 적용합니다. 안전 평가는 원청 안전팀·감시단·본사점검·자체점검의 지적 이력을 기준으로 감점하는 방식으로 운영하며, 별도 문제가 없는 경우 감점하지 않습니다.",
            "안전평가의 객관성을 확보하기 위해 원청/도급사 안전팀, 감시단, 본사 안전점검, 현장 자체점검에서 확인된 지적사항을 제재 또는 감점 근거로 등록합니다. 이를 통해 현장 내부 평가자의 주관성을 줄이고, 실제 안전수칙 위반 이력에 기반하여 안전점수를 반영합니다.",
        ],
        "05_evaluate_safety.png",
    )

    screens = [
        ("3-1. 접속 및 로그인", "01_login.png", [
            "BESMA 로그인 화면에서 평가자 계정(팀장/소장/본사)으로 접속",
            "현장 평가자: /site/functional-eval · 본사: /hq-safe/functional-eval",
        ], ["주의: 평가자 계정은 ERP 출역 반영 후 자동 생성(별칭-이름 / 주민앞6자리)"]),
        ("3-2. 최초 로그인 — 평가책임 동의서", "02_consent_full.png", [
            "기능인제 화면 최초 진입 시 동의서 모달 표시(1회)",
            "동의문을 끝까지 스크롤 → 「위 내용을 확인하였으며 동의합니다」 체크 → 캔버스 서명",
            "「동의 및 서명」 클릭 시 PDF·서명 PNG 서버 저장",
        ], ["완료 조건: consent/status.required = false", "스크롤 게이트: 끝까지 확인 전 체크·서명·제출 비활성화"], "02_consent_bottom.png"),
        ("3-3. 근로자 명단(등급현황)", "03_roster_team.png", [
            "담당 근로자 목록 · 2-1/2-2 등급 · 안전·제재 상태 · 비고(포상/감점) 확인",
            "「평가」 버튼으로 2-1/2-2 입력 · 「포상」「제재」 버튼으로 근거 등록",
        ], ["완료 조건: 담당 전원 평가 완료 + (필요 시) 포상·제재 근거 등록"]),
        ("3-4. 2-1 기능 평가 입력", "04_evaluate_functional.png", [
            "문항별 우수/보통/부족/문제 선택 후 저장",
            "미평가자는 사이드바 「미평가」 필터로 확인",
        ], ["완료 조건: FUNCTIONAL 전 문항 입력(is_complete)"]),
        ("3-5. 2-2 안전 평가 입력", "05_evaluate_safety.png", [
            "안전 8개 문항 등급 입력 · 하단 「위반·제재」 인라인 등록 가능",
            "「문제」 항목이 있으면 제재 등록 유도(완료 차단은 아님)",
        ], ["완료 조건: SAFETY 전 문항 입력"]),
        ("3-6. 고객사 포상 등록", "06_reward_modal.png", [
            "명단 「포상」 → 사진 첨부 → 제출(본사 승인 대기)",
            "승인 후 비고 「고객사포상(+5)」 · 안전 등급 +5점 반영",
        ], ["주의: 제출 후 회수·변경 불가 · 소장은 팀원 포상 직접 등록 불가"]),
        ("3-7. 제재 등록", "07_sanction_modal.png", [
            "명단 「제재」 또는 2-2 탭 → 위반 선택 · 근거(코멘트/사진) · 서명",
            "첫 제재: 해당 안전항목 「문제」 자동 반영(감점 0)",
            "같은 위반 재발: -5점 감점(평가표 유지) · 안전 등급에 감점 합산 반영",
        ], ["평가 완료 후에도 제재 등록 가능 · 제재 누락으로 평가완료 서명 자체는 차단하지 않음"]),
    ]
    for item in screens:
        if len(item) == 5:
            title, img, bullets, notes, img2 = item
            _slide_screen(prs, title, img, bullets, notes, img2)
        else:
            title, img, bullets, notes = item
            _slide_screen(prs, title, img, bullets, notes)

    # 포상/제재 점수 로직
    _slide_screen(
        prs,
        "4. 포상·제재 — 점수·등급 반영 로직",
        "07_sanction_modal.png",
        [
            "2-1 기능/품질: 문항 등급 합산 → S/A/B/C · 기본 B · S 20% 권장(기능만)",
            "2-2 안전: 감점형 — 지적·제재 이력 없으면 감점 없음 · S/C 비율 제한 없음",
            "첫 제재: 해당 안전항목 「문제」 반영(감점 0) · 재발 -5점",
            "포상 승인: bonus_points 가점 · 안전 등급 재산출",
            "감점 근거: 원청/도급 안전팀·감시단·본사점검·자체점검·보호구·고소·작업허가·정리정돈·반복지적·사고",
        ],
        ["제재 등록 권장(유도) · 미등록 시에도 평가완료 서명 가능 · 안전 S 20% 제한 없음"],
    )

    # 평가완료 서명
    _slide_screen(
        prs,
        "5. 평가완료보고서 서명",
        "12_team_signoff_modal.png" if (SHOT_DIR / "12_team_signoff_modal.png").is_file() else "09_manager_approval_panel.png",
        [
            "팀장: 담당 팀원 전원 2-1+2-2 완료 → 「평가완료보고서 서명」(TEAM_LEADER stage PDF)",
            "소장: (11명↑) 팀장 전원 서명 확인 → 팀장보고서 승인 → 「평가완료보고서 제출」(SITE stage)",
            "서명 시 캔버스 입력 → PDF 자동 생성 · /functional-eval/signatures/{id}/document 다운로드",
            "서명 완료 후 evaluation_editable=false → 점수·포상·제재 수정 차단",
            "등급 인플레이션 방지: 기능/품질(2-1)만 S 20% 초과 시 서명 전 사유 · 안전(2-2) 제한 없음",
        ],
        [
            "완료 조건: 출역 대상 전원 평가 완료 + 필수 서명 단계 완료",
            "기능/품질 평가는 소장·팀장이 작업숙련도와 품질을 기준으로 평가하고, S등급 남발 방지를 위해 기본 B등급 및 S등급 권장 상한을 적용합니다. 안전 평가는 원청 안전팀·감시단·본사점검·자체점검의 지적 이력을 기준으로 감점하는 방식으로 운영하며, 별도 문제가 없는 경우 감점하지 않습니다.",
        ],
    )

    # 결재
    _slide_screen(
        prs,
        "6-1. 소장 승인",
        "08_manager_roster.png",
        [
            "계정: 대우청라-박명식(소장) / 현장별 {별칭}-{소장명}",
            "화면: /site/functional-eval — 직영 평가 + 팀장 보고서 승인 패널",
            "확인: 팀원 등급표·포상/제재 근거 → 팀장보고서 「승인 서명」→ 전원 완료 후 「제출 및 서명」",
        ],
    )
    _slide_screen(
        prs,
        "6-2. 안전보건실 담당 검토·승인",
        "10_hq_dashboard.png",
        [
            "계정: 안전보건-정상익 차장 /hq-safe/functional-eval",
            "소장 제출(SITE_APPROVED) 현장 → 담당 검토 코멘트 입력",
            "「담당 일괄 검토·승인」→ HQ_OFFICER stage PDF · 상태 HQ_OFFICER_APPROVED",
        ],
    )
    _slide_screen(
        prs,
        "6-3. 안전보건실장 최종 승인",
        "10_hq_dashboard.png",
        [
            "계정: 안전보건-조동문 전무 /hq-safe/functional-eval",
            "담당 승인 완료 현장 → 실장 코멘트 입력",
            "「실장 일괄 최종승인」→ HQ stage PDF · 상태 HQ_APPROVED",
        ],
    )
    _slide_screen(
        prs,
        "6-4. 대표이사 최종 승인",
        "11_ceo_dashboard.png",
        [
            "계정: 부현대표-김홍수 /hq-safe/functional-eval",
            "HQ 승인(HQ_APPROVED) 현장 최종 확인 → 「최종 승인 및 서명」",
            "CEO stage PDF · 상태 CEO_APPROVED · 현장 평가 확정",
        ],
    )

    # 최종 보고서
    _slide_screen(
        prs,
        "7. 최종 평가보고서 PDF",
        "10_hq_dashboard.png",
        [
            "제목: 기능인인정제 평가 보고서 (갑지/상세/본사·대표 단계 공통)",
            "포함: 평가기간 · 현장명 · 피평가자 등급표 · 서명 이미지 · 결재란",
            "다운로드: GET /functional-eval/signatures/mine · signatures/{id}/document",
            "동의서 PDF: GET /functional-eval/consent/document",
            "등급표 엑셀: /my-site/export/site-grade-workbook",
        ],
    )

    # 추가평가
    _slide_screen(
        prs,
        "8. 신규입사·추가평가 처리",
        None,
        [
            "핵심 원칙: 「기존 서명본은 수정하지 않고, 신규 대상자만 별도 추가평가로 생성·평가·서명」",
            "evaluation_batch=0(최초) / 1+(추가평가 N차) — 신규 worker 생성 시 자동 부여",
            "API: POST /functional-eval/my-site/supplemental-signoff (기존 CEO_APPROVED 유지)",
            "신규입사자 추가평가는 기존 평가보고서를 수정하지 않고 별도 배치로 분리하는 구조가 반영되어 있으며, 전용 화면은 후속 보완 예정입니다.",
        ],
    )

    # 검증 결과
    _slide_table(
        prs,
        "9. 시스템 검증 결과 (2026-06-15)",
        ["검증 항목", "결과", "비고"],
        [
            ["최초 로그인 동의서 표시", "통과", "FeConsentGate · API 테스트"],
            ["동의서 하단 스크롤 전 서명 제한", "통과", "FeSignatureModal requireConsentScroll"],
            ["동의서 서명 후 재로그인 시 미표시", "통과", "test_consent_submit_once"],
            ["근로자 2-1/2-2 평가 입력", "통과", "E2E 시뮬레이션 + 화면 캡처"],
            ["포상/제재 등록", "통과", "Form+서명 API · UI 캡처"],
            ["점수/등급 반영(감점·가점)", "통과", "test_sanction_register_without_safety_bottom"],
            ["평가보고서 PDF 제목 통일", "통과", "기능인인정제 평가 보고서"],
            ["등급 인플레이션 방지(기능/품질 S 20%)", "통과", "2-1만 적용 · 안전 2-2 제한 없음"],
            ["평가보고서 서명", "통과", "TEAM/SITE/HQ/CEO stage PDF"],
            ["서명 후 수정 차단", "통과", "test_signature_lock_after_team_signoff"],
            ["소장 승인", "통과", "E2E simulation step 2"],
            ["안전보건실 담당·실장 2단 승인", "통과", "HQ officer → director"],
            ["대표이사 최종 승인", "통과", "E2E simulation step 4"],
            ["PDF 다운로드", "통과", "consent/signatures/document API"],
            ["추가평가 분리 처리", "부분", "API·batch 구현 / 전용 UI 후속"],
        ],
    )

    _slide_screen(
        prs,
        "10. 시연·운영 전 안내",
        None,
        [
            "금일 점검 과정에서 동의·서명·현장승인 레코드는 초기화되었습니다.",
            "따라서 시연 및 실제 운영 전 각 사용자는 최초 로그인 동의서를 다시 서명해야 합니다.",
            "평가 점수, 출역, 근로자, 삼성인정제 관련 데이터는 유지되었습니다.",
        ],
    )

    # 요약
    _slide_screen(
        prs,
        "11. 전체 운영 흐름 요약",
        "03_roster_team.png",
        [
            "출역 ERP → 평가자·피평가자 자동 명단 → 동의(1회) → 2-1/2-2 평가",
            "포상(가점)·제재(감점/문제반영) → 팀장·소장 서명 → 담당→실장→대표 3단 본사 승인 → PDF 영구보관",
            "서명 후 수정 불가 · 추가 인원은 추가평가로 분리",
            "대표이사 보고 시: 본 자료 + E2E 리포트(docs/reports/functional-eval-e2e/) 참조",
        ],
    )

    OUT_PPT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PPT))
    return OUT_PPT


if __name__ == "__main__":
    path = build()
    print(f"Wrote {path}")
